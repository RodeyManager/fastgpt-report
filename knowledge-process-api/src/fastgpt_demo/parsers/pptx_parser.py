"""PPTX parser — extracts text from all slides and shapes."""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from ._types import ParseResult


def parse(buffer: bytes, **kwargs) -> ParseResult:
    """Parse a PPTX buffer, extracting text from every slide."""

    prs = Presentation(BytesIO(buffer))

    slide_texts: list[str] = []

    for slide in prs.slides:
        shape_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        shape_texts.append(text)
        if shape_texts:
            slide_texts.append("\n".join(shape_texts))

    raw_text = "\n".join(slide_texts)

    return ParseResult(
        raw_text=raw_text,
        format_text=raw_text,
        html_preview=f"<pre>{raw_text}</pre>",
        image_list=[],
    )
